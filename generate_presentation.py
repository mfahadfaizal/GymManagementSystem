#!/usr/bin/env python3
"""
Gym Management System - Professional PowerPoint Presentation Generator
This script creates a comprehensive presentation with all necessary diagrams
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.drawing.line import LineFormat
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

def create_title_slide(prs):
    """Create the title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Gym Management System"
    subtitle.text = "Comprehensive System Overview & Architecture\nProfessional Presentation"
    
    # Style the title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

def create_agenda_slide(prs):
    """Create agenda slide"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Presentation Agenda"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    agenda_items = [
        "System Overview & Business Value",
        "Technical Architecture",
        "Database Design & Schema",
        "User Roles & Access Control",
        "Core Features & Functionality",
        "API Endpoints & Integration",
        "Security & Authentication",
        "Deployment & Scalability",
        "Future Enhancements",
        "Q&A & Discussion"
    ]
    
    for item in agenda_items:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.space_after = Pt(6)

def create_system_overview_slide(prs):
    """Create system overview slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "System Overview & Business Value"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    overview_text = [
        "🏋️ Full-Stack Gym Management Solution",
        "🎯 Comprehensive Role-Based Access Control",
        "💼 Streamlined Operations Management",
        "📊 Real-Time Analytics & Reporting",
        "🔐 Enterprise-Grade Security",
        "📱 Responsive Web Application",
        "🗄️ Robust Database Architecture",
        "🔄 Seamless API Integration"
    ]
    
    for text in overview_text:
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.space_after = Pt(8)

def create_architecture_diagram_slide(prs):
    """Create architecture diagram slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "System Architecture"
    
    # Create a simple architecture diagram using shapes
    left = Inches(1)
    top = Inches(2)
    width = Inches(2)
    height = Inches(1.5)
    
    # Frontend Layer
    frontend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    frontend.text = "Frontend\nReact.js\nBootstrap 5"
    frontend.fill.solid()
    frontend.fill.fore_color.rgb = (70, 130, 180)  # Steel Blue
    
    # Backend Layer
    backend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + 3, top, width, height)
    backend.text = "Backend\nSpring Boot\nJava 17"
    backend.fill.solid()
    backend.fill.fore_color.rgb = (46, 139, 87)  # Sea Green
    
    # Database Layer
    database = slide.shapes.add_shape(MSO_SHAPE.CYLINDER, left + 1.5, top + 2.5, width, height)
    database.text = "Database\nMySQL 8.0\nJPA/Hibernate"
    database.fill.solid()
    database.fill.fore_color.rgb = (255, 140, 0)  # Dark Orange
    
    # Add arrows
    arrow1 = slide.shapes.add_connector(1, left + 2, top + 0.75, left + 3, top + 0.75)
    arrow2 = slide.shapes.add_connector(1, left + 3, top + 1.5, left + 2.5, top + 2.5)

def create_database_schema_slide(prs):
    """Create database schema slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Database Schema Design"
    
    # Create entity boxes
    entities = [
        ("Users", 1, 2, "Role-based access\nAuthentication"),
        ("Memberships", 3, 2, "Subscription plans\nStatus tracking"),
        ("Equipment", 5, 2, "Inventory management\nMaintenance tracking"),
        ("Gym Classes", 1, 4, "Class scheduling\nCapacity management"),
        ("Training Sessions", 3, 4, "Personal training\nBooking system"),
        ("Payments", 5, 4, "Transaction records\nMultiple methods"),
        ("Class Registrations", 3, 6, "Attendance tracking\nEnrollment management")
    ]
    
    for entity_name, x, y, description in entities:
        left = Inches(x)
        top = Inches(y)
        width = Inches(2)
        height = Inches(1.2)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.text = f"{entity_name}\n{description}"
        box.fill.solid()
        box.fill.fore_color.rgb = (135, 206, 235)  # Sky Blue

def create_user_roles_slide(prs):
    """Create user roles slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "User Roles & Access Control"
    
    # Create role boxes
    roles = [
        ("ADMIN", "Full system access\nUser management\nSystem configuration", 1, 2),
        ("STAFF", "Membership management\nEquipment management\nPayment processing", 3, 2),
        ("TRAINER", "Training sessions\nClass teaching\nProgress tracking", 5, 2),
        ("MEMBER", "Profile management\nSession booking\nClass registration", 3, 4)
    ]
    
    for role_name, description, x, y in roles:
        left = Inches(x)
        top = Inches(y)
        width = Inches(2.5)
        height = Inches(1.5)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.text = f"{role_name}\n{description}"
        box.fill.solid()
        
        # Different colors for different roles
        if role_name == "ADMIN":
            box.fill.fore_color.rgb = (220, 20, 60)  # Crimson
        elif role_name == "STAFF":
            box.fill.fore_color.rgb = (255, 165, 0)  # Orange
        elif role_name == "TRAINER":
            box.fill.fore_color.rgb = (34, 139, 34)  # Forest Green
        else:  # MEMBER
            box.fill.fore_color.rgb = (70, 130, 180)  # Steel Blue

def create_features_slide(prs):
    """Create features slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Core Features & Functionality"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    features = [
        "🔐 JWT Authentication & Authorization",
        "👥 Multi-Role User Management",
        "💳 Membership & Subscription Management",
        "🏋️ Equipment Inventory & Maintenance",
        "📚 Class Scheduling & Registration",
        "👨‍🏫 Personal Training Session Booking",
        "💰 Payment Processing & Tracking",
        "📊 Real-Time Analytics & Reporting",
        "🔍 Advanced Search & Filtering",
        "📱 Responsive Web Interface"
    ]
    
    for feature in features:
        p = text_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(16)
        p.space_after = Pt(6)

def create_api_endpoints_slide(prs):
    """Create API endpoints slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "API Endpoints & Integration"
    
    # Create API categories
    categories = [
        ("Authentication", ["POST /api/auth/signin", "POST /api/auth/signup"], 1, 2),
        ("Users", ["GET /api/users", "PUT /api/users/{id}", "DELETE /api/users/{id}"], 3, 2),
        ("Memberships", ["GET /api/memberships", "POST /api/memberships", "PUT /api/memberships/{id}"], 5, 2),
        ("Equipment", ["GET /api/equipment", "POST /api/equipment", "PUT /api/equipment/{id}"], 1, 4),
        ("Classes", ["GET /api/gym-classes", "POST /api/gym-classes", "PUT /api/gym-classes/{id}"], 3, 4),
        ("Sessions", ["GET /api/training-sessions", "POST /api/training-sessions/book"], 5, 4)
    ]
    
    for category, endpoints, x, y in categories:
        left = Inches(x)
        top = Inches(y)
        width = Inches(2.5)
        height = Inches(1.8)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.text = f"{category}\n\n" + "\n".join(endpoints)
        box.fill.solid()
        box.fill.fore_color.rgb = (240, 248, 255)  # Alice Blue

def create_security_slide(prs):
    """Create security slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Security & Authentication"
    
    # Create security components
    components = [
        ("JWT Tokens", "Secure token-based\nauthentication", 1, 2),
        ("BCrypt", "Password encryption\n& hashing", 3, 2),
        ("Spring Security", "Role-based access\ncontrol", 5, 2),
        ("CORS", "Cross-origin\nresource sharing", 1, 4),
        ("Input Validation", "Data sanitization\n& validation", 3, 4),
        ("HTTPS", "Secure communication\nprotocols", 5, 4)
    ]
    
    for component, description, x, y in components:
        left = Inches(x)
        top = Inches(y)
        width = Inches(2.5)
        height = Inches(1.5)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.text = f"{component}\n{description}"
        box.fill.solid()
        box.fill.fore_color.rgb = (255, 215, 0)  # Gold

def create_deployment_slide(prs):
    """Create deployment slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Deployment & Scalability"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    deployment_info = [
        "🚀 Spring Boot Application Server",
        "🗄️ MySQL Database with Connection Pooling",
        "🌐 React.js Frontend with Bootstrap",
        "🔧 Maven Build & Dependency Management",
        "📦 Docker Containerization Ready",
        "☁️ Cloud Deployment Compatible",
        "📊 Monitoring & Logging Integration",
        "🔄 CI/CD Pipeline Support"
    ]
    
    for info in deployment_info:
        p = text_frame.add_paragraph()
        p.text = info
        p.font.size = Pt(16)
        p.space_after = Pt(6)

def create_future_enhancements_slide(prs):
    """Create future enhancements slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Future Enhancements & Roadmap"
    
    # Create enhancement categories
    enhancements = [
        ("Mobile App", "iOS & Android\napplications", 1, 2),
        ("AI Integration", "Smart scheduling\n& recommendations", 3, 2),
        ("Payment Gateway", "Stripe, PayPal\nintegration", 5, 2),
        ("Analytics", "Advanced reporting\n& insights", 1, 4),
        ("Notifications", "Email & SMS\nalerts", 3, 4),
        ("IoT Integration", "Smart equipment\nmonitoring", 5, 4)
    ]
    
    for enhancement, description, x, y in enhancements:
        left = Inches(x)
        top = Inches(y)
        width = Inches(2.5)
        height = Inches(1.5)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.text = f"{enhancement}\n{description}"
        box.fill.solid()
        box.fill.fore_color.rgb = (138, 43, 226)  # Blue Violet

def create_qa_slide(prs):
    """Create Q&A slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Questions & Discussion"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    qa_text = [
        "❓ Technical Questions",
        "💡 Feature Requests",
        "🔧 Implementation Details",
        "📊 Performance Considerations",
        "🚀 Deployment Strategies",
        "🔄 Integration Options",
        "📱 User Experience",
        "🔐 Security Concerns"
    ]
    
    for text in qa_text:
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(18)
        p.space_after = Pt(8)

def create_contact_slide(prs):
    """Create contact slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Thank You!"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    contact_text = [
        "🎯 Gym Management System",
        "📧 Contact: development@example.com",
        "🌐 Website: www.example.com",
        "📱 Phone: +1 (555) 123-4567",
        "",
        "Ready for implementation and deployment!",
        "",
        "Any additional questions?"
    ]
    
    for text in contact_text:
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.space_after = Pt(6)
        
        if "Ready for implementation" in text:
            p.font.bold = True
            p.font.size = Pt(18)

def main():
    """Main function to create the presentation"""
    print("Creating Gym Management System PowerPoint Presentation...")
    
    # Create presentation
    prs = Presentation()
    
    # Add slides
    create_title_slide(prs)
    create_agenda_slide(prs)
    create_system_overview_slide(prs)
    create_architecture_diagram_slide(prs)
    create_database_schema_slide(prs)
    create_user_roles_slide(prs)
    create_features_slide(prs)
    create_api_endpoints_slide(prs)
    create_security_slide(prs)
    create_deployment_slide(prs)
    create_future_enhancements_slide(prs)
    create_qa_slide(prs)
    create_contact_slide(prs)
    
    # Save presentation
    output_file = "Gym_Management_System_Presentation.pptx"
    prs.save(output_file)
    
    print(f"✅ Presentation created successfully: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("🎯 Ready for professional presentation!")

if __name__ == "__main__":
    main() 